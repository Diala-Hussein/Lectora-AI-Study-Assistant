from pptx import Presentation
from typing import Optional
import io

class PPTXExtractor:
    @staticmethod
    def extract(uploaded_file) -> Optional[str]:
        """Extract text from PPTX file"""
        text = ""
        try:
            # Read uploaded file
            prs = Presentation(io.BytesIO(uploaded_file.read()))
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = ""
                
                # Extract text from all shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for paragraph in shape.text_frame.paragraphs:
                            slide_text += paragraph.text + " "
                    elif hasattr(shape, "text"):
                        slide_text += shape.text + " "
                
                # Only add non-empty slides
                if slide_text.strip():
                    text += f"\n--- Slide {slide_num} ---\n{slide_text.strip()}\n"
            
            return text.strip()
            
        except Exception as e:
            print(f"PPTX extraction error: {e}")
            return ""